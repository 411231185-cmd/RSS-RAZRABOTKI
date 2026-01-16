from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

print("🔄 Начинаю конвертацию...")

# Читаем markdown файл
with open('Prezentacija_Zvezdochka77.md', 'r', encoding='utf-8') as f:
    content = f.read()

# Создаём презентацию
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Разделяем контент на слайды
slides_raw = content.split('<div style="page-break-after: always;"></div>')

slide_count = 0

for slide_text in slides_raw:
    if not slide_text.strip() or slide_text.strip() == '---':
        continue
    
    # Создаём новый слайд (пустой макет)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Парсим контент слайда
    lines = [line for line in slide_text.split('\n') if line.strip()]
    
    # Ищем заголовок
    title_text = ""
    content_lines = []
    
    for line in lines:
        clean_line = line.strip()
        if clean_line.startswith('# ') and not title_text:
            title_text = clean_line.lstrip('#').strip()
        elif clean_line.startswith('## '):
            if not title_text:
                title_text = clean_line.lstrip('#').strip()
            else:
                content_lines.append(clean_line.lstrip('#').strip())
        elif clean_line.startswith('### '):
            content_lines.append('• ' + clean_line.lstrip('#').strip())
        elif clean_line and clean_line != '---' and not clean_line.startswith('**[МЕСТО ДЛЯ ФОТО'):
            content_lines.append(clean_line)
    
    # Добавляем заголовок на слайд
    if title_text:
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title_text
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)  # Синий цвет
        p.alignment = PP_ALIGN.CENTER
    
    # Добавляем содержимое
    if content_lines:
        left = Inches(0.7)
        top = Inches(1.5)
        width = Inches(8.6)
        height = Inches(3.8)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        # Добавляем первые строки (чтобы не переполнить слайд)
        max_lines = 20
        for line in content_lines[:max_lines]:
            p = text_frame.add_paragraph()
            p.text = line.replace('**', '').replace('|', ' ')
            p.font.size = Pt(14)
            p.space_after = Pt(6)
            
            if line.startswith('•'):
                p.level = 1
    
    slide_count += 1
    print(f"✅ Слайд {slide_count}: {title_text[:50]}...")

# Сохраняем презентацию
output_file = 'Prezentacija_Zvezdochka77.pptx'
prs.save(output_file)

print(f"\n🎉 Готово! Создано слайдов: {slide_count}")
print(f"📁 Файл сохранён: {output_file}")
print("\n💡 Откройте файл в PowerPoint и добавьте фотографии!")
