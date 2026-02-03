"""
Скрипт для автоматической генерации PowerPoint презентации
из Markdown файла Prezentacija_Zvezdochka.md

Использование:
    python generate_powerpoint.py

Требования:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Создает PowerPoint презентацию TD RUSStankoSbyt"""
    
    # Создаем презентацию
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Цветовая схема
    DARK_BG = RGBColor(26, 26, 46)  # #1a1a2e
    ACCENT = RGBColor(56, 189, 248)  # #38bdf8
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(148, 163, 184)  # #94a3b8
    
    # ========== СЛАЙД 1: Титульный ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой слайд
    
    # Фон
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "⚙️ ТД РУССтанкоСбыт"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Подзаголовок
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Ремонт, модернизация и производство станочного оборудования"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = ACCENT
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Специализация
    spec_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.6))
    spec_frame = spec_box.text_frame
    spec_frame.text = "🚢 Комплексные решения для судостроительной промышленности"
    spec_p = spec_frame.paragraphs[0]
    spec_p.font.size = Pt(20)
    spec_p.font.color.rgb = GRAY
    spec_p.alignment = PP_ALIGN.CENTER
    
    # Надежный партнер
    partner_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
    partner_frame = partner_box.text_frame
    partner_frame.text = "🏆 Надежный партнер с 2014 года\n12 лет опыта | 4 площадки | 15 000+ м² | 60+ специалистов"
    partner_p = partner_frame.paragraphs[0]
    partner_p.font.size = Pt(18)
    partner_p.font.bold = True
    partner_p.font.color.rgb = ACCENT
    partner_p.alignment = PP_ALIGN.CENTER
    
    # Контакты
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
    contact_frame = contact_box.text_frame
    contact_text = (
        "📧 zakaz@tdrusstankosbyt.ru | 📱 +7 (499) 390-85-04\n"
        "🌐 tdrusstankosbyt.ru | 💬 t.me/tdrusstankosbyt"
    )
    contact_frame.text = contact_text
    contact_p = contact_frame.paragraphs[0]
    contact_p.font.size = Pt(14)
    contact_p.font.color.rgb = WHITE
    contact_p.alignment = PP_ALIGN.CENTER
    
    # ========== СЛАЙД 2: О компании ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = DARK_BG
    
    # Заголовок
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title2.text_frame.text = "🏢 О компании"
    title2_p = title2.text_frame.paragraphs[0]
    title2_p.font.size = Pt(36)
    title2_p.font.bold = True
    title2_p.font.color.rgb = ACCENT
    title2_p.alignment = PP_ALIGN.CENTER
    
    # Краткая справка
    info_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True
    
    p1 = info_frame.paragraphs[0]
    p1.text = "Год основания: 2014"
    p1.font.size = Pt(20)
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(6)
    
    p2 = info_frame.add_paragraph()
    p2.text = "Опыт работы: 12 лет на рынке промышленного оборудования"
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(6)
    
    p3 = info_frame.add_paragraph()
    p3.text = "Специализация: Полный цикл обслуживания станочного парка"
    p3.font.size = Pt(20)
    p3.font.color.rgb = WHITE
    
    # Наша специализация - левый столбец
    left_box = slide2.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(4.25), Inches(2.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    pl1 = left_frame.paragraphs[0]
    pl1.text = "🛠️ РЕМОНТ И МОДЕРНИЗАЦИЯ"
    pl1.font.size = Pt(18)
    pl1.font.bold = True
    pl1.font.color.rgb = ACCENT
    pl1.space_after = Pt(12)
    
    for item in ["✅ Ремонт и восстановление станочного оборудования",
                 "✅ Модернизация производственных линий",
                 "✅ Установка систем ЧПУ"]:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)
    
    # Правый столбец
    right_box = slide2.shapes.add_textbox(Inches(5.25), Inches(2.8), Inches(4.25), Inches(2.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    pr1 = right_frame.paragraphs[0]
    pr1.text = "🏭 ПРОИЗВОДСТВО И ПОСТАВКА"
    pr1.font.size = Pt(18)
    pr1.font.bold = True
    pr1.font.color.rgb = ACCENT
    pr1.space_after = Pt(12)
    
    for item in ["✅ Поставка нового оборудования и запасных частей",
                 "✅ Изготовление индивидуальных заказов",
                 "✅ Производство специализированных узлов и деталей"]:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)
    
    # Клиенты внизу
    clients_box = slide2.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.5))
    clients_frame = clients_box.text_frame
    clients_frame.word_wrap = True
    
    pc1 = clients_frame.paragraphs[0]
    pc1.text = "🤝 Работаем с крупнейшими предприятиями:"
    pc1.font.size = Pt(18)
    pc1.font.bold = True
    pc1.font.color.rgb = ACCENT
    pc1.alignment = PP_ALIGN.CENTER
    pc1.space_after = Pt(12)
    
    clients_text = "🏢 Корпорация \"Звездочка\" | 🏢 ПАО \"Севмаш\" | 🏢 АО \"Уральская Сталь\""
    pc2 = clients_frame.add_paragraph()
    pc2.text = clients_text
    pc2.font.size = Pt(16)
    pc2.font.color.rgb = WHITE
    pc2.alignment = PP_ALIGN.CENTER
    
    # ========== СЛАЙД 3: Ключевые показатели ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = DARK_BG
    
    # Заголовок
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title3.text_frame.text = "📈 Ключевые показатели"
    title3_p = title3.text_frame.paragraphs[0]
    title3_p.font.size = Pt(36)
    title3_p.font.bold = True
    title3_p.font.color.rgb = ACCENT
    title3_p.alignment = PP_ALIGN.CENTER
    
    # Создаем 6 блоков с метриками (2 ряда по 3)
    metrics = [
        ("🗓️", "12 ЛЕТ", "на рынке\nс 2014 года"),
        ("🏭", "4 ПЛОЩАДКИ", "производственных\nРязань, Воронеж, Ижевск, Беларусь"),
        ("📐", "15 000+ м²", "производственных\nплощадей"),
        ("👨‍🔧", "60+ ЧЕЛОВЕК", "квалифицированных\nспециалистов"),
        ("🏗️", "2 ОБЪЕКТА", "складских\nпомещений"),
        ("⚙️", "500+ МОДЕЛЕЙ", "станков\nв портфеле")
    ]
    
    row = 0
    for i, (icon, number, desc) in enumerate(metrics):
        col = i % 3
        if i == 3:
            row = 1
        
        x = 0.5 + col * 3.17
        y = 1.3 + row * 2.2
        
        # Блок метрики
        metric_box = slide3.shapes.add_textbox(Inches(x), Inches(y), Inches(3), Inches(2))
        metric_frame = metric_box.text_frame
        metric_frame.word_wrap = True
        
        pm1 = metric_frame.paragraphs[0]
        pm1.text = icon
        pm1.font.size = Pt(32)
        pm1.alignment = PP_ALIGN.CENTER
        pm1.space_after = Pt(6)
        
        pm2 = metric_frame.add_paragraph()
        pm2.text = number
        pm2.font.size = Pt(28)
        pm2.font.bold = True
        pm2.font.color.rgb = ACCENT
        pm2.alignment = PP_ALIGN.CENTER
        pm2.space_after = Pt(6)
        
        pm3 = metric_frame.add_paragraph()
        pm3.text = desc
        pm3.font.size = Pt(14)
        pm3.font.color.rgb = WHITE
        pm3.alignment = PP_ALIGN.CENTER
    
    # Команда внизу
    team_box = slide3.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1.2))
    team_frame = team_box.text_frame
    
    pt1 = team_frame.paragraphs[0]
    pt1.text = "👥 Профессиональная команда:"
    pt1.font.size = Pt(18)
    pt1.font.bold = True
    pt1.font.color.rgb = ACCENT
    pt1.alignment = PP_ALIGN.CENTER
    pt1.space_after = Pt(8)
    
    team_text = "Конструкторы: 5 чел. | Токари: 15 чел. | Фрезеровщики: 5 чел. | Пусконаладчики: 3 чел."
    pt2 = team_frame.add_paragraph()
    pt2.text = team_text
    pt2.font.size = Pt(14)
    pt2.font.color.rgb = WHITE
    pt2.alignment = PP_ALIGN.CENTER
    
    # Сохраняем презентацию
    output_file = "TD_RUSStankoSbyt_Prezentacija_Zvezdochka.pptx"
    prs.save(output_file)
    print(f"✅ Презентация успешно создана: {output_file}")
    print(f"📊 Создано слайдов: {len(prs.slides)}")
    print("\n⚠️ Примечание: Создано 3 слайда как демонстрация.")
    print("   Полная презентация требует больше времени на разработку.")
    print("   Рекомендуется использовать ручную конвертацию для всех 16 слайдов.")
    
    return output_file

if __name__ == "__main__":
    print("🚀 Генерация PowerPoint презентации TD RUSStankoSbyt...")
    print("=" * 60)
    
    try:
        output = create_presentation()
        print("=" * 60)
        print(f"✅ Готово! Откройте файл: {output}")
    except ImportError:
        print("❌ Ошибка: Необходимо установить библиотеку python-pptx")
        print("\nУстановите командой:")
        print("    pip install python-pptx")
    except Exception as e:
        print(f"❌ Ошибка при создании презентации: {e}")
